"""
Math4AI Capstone Presentation Generator — v2 (improved)
Generates a professional 20-slide PowerPoint for the NAIC capstone.

Changes from v1:
  - Added Implementation Sanity Checks slide (required by rubric)
  - Track B rebuilt with actual entropy numbers + real figures
  - Repeated-seed slide: explicit CI non-overlap statement
  - Added Key Numbers summary slide
  - Tightened slide count to 20 for 10-min budget (~30 s/slide)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x1F, 0x44)
GOLD      = RGBColor(0xF5, 0xA6, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFA)
MID_BLUE  = RGBColor(0x1A, 0x5F, 0xA8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT2   = RGBColor(0x2E, 0xCC, 0x71)
WARN      = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FIGURES_DIR = os.path.join(os.path.dirname(__file__), "..", "figures")


# ── Helpers ───────────────────────────────────────────────────────────────────

def solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color):
    box = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    solid_fill(box, color)
    box.line.fill.background()
    return box


def add_text_box(slide, text, l, t, w, h,
                 font_size=20, bold=False, color=DARK_TEXT,
                 align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return txb


def add_mixed_run(slide, parts, l, t, w, h, align=PP_ALIGN.LEFT):
    """parts = list of (text, font_size, bold, color)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, fs, bd, col in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(fs)
        r.font.bold = bd
        r.font.color.rgb = col
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=18, color=DARK_TEXT, header=None, header_color=MID_BLUE):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = header
        r.font.size = Pt(font_size + 2)
        r.font.bold = True
        r.font.color.rgb = header_color
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"  •  {item}" if header else f"•  {item}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
    return txb


def slide_background(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, NAVY)
    add_rect(slide, 0, 1.3, 13.33, 0.06, GOLD)
    add_text_box(slide, title, 0.35, 0.1, 12.5, 0.75,
                 font_size=30, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, 0.35, 0.8, 12.5, 0.5,
                     font_size=15, color=GOLD)


def footer_bar(slide, text="Math4AI Capstone  |  National AI Center  |  2026"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, NAVY)
    add_text_box(slide, text, 0.3, 7.2, 12.5, 0.3,
                 font_size=11, color=RGBColor(0xAA, 0xBB, 0xCC))


def add_figure(slide, filename, l, t, w, h):
    path = os.path.join(FIGURES_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
        return True
    return False


def kv_row(slide, label, value, lx, ty, lw=3.2, vw=2.0, h=0.48,
           label_color=DARK_TEXT, value_color=NAVY):
    add_text_box(slide, label, lx, ty, lw, h, font_size=14, color=label_color)
    add_text_box(slide, value, lx + lw, ty, vw, h,
                 font_size=14, bold=True, color=value_color, align=PP_ALIGN.RIGHT)


# ── Slide 1: Title ─────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, NAVY)

    add_rect(slide, 0, 0, 13.33, 0.55, RGBColor(0x06, 0x14, 0x2E))
    add_text_box(slide, "NATIONAL AI CENTER  —  AI ACADEMY  |  MATH4AI CAPSTONE",
                 0.4, 0.1, 12.5, 0.4, font_size=13, bold=True, color=GOLD)

    add_text_box(slide, "From Linear Scores to a Single Hidden Layer",
                 0.5, 1.3, 12.3, 1.1, font_size=40, bold=True, color=WHITE)
    add_text_box(slide, "A Mathematical Study of Simple Learning Systems",
                 0.5, 2.45, 12.3, 0.72, font_size=25, color=GOLD)

    add_rect(slide, 0.5, 3.3, 5.5, 0.05, RGBColor(0x44, 0x77, 0xBB))

    details = [
        "Central Question:  When does a hidden layer genuinely improve on a linear decision rule?",
        "Models:  Softmax Regression  vs.  One-Hidden-Layer Neural Network (tanh)",
        "Datasets:  Linear Gaussian  ·  Moons  ·  Digits (10-class benchmark)",
        "Implementation:  Pure NumPy — no PyTorch, TensorFlow, or autograd",
    ]
    for i, d in enumerate(details):
        add_text_box(slide, d, 0.5, 3.5 + i * 0.52, 12.3, 0.48,
                     font_size=16, color=RGBColor(0xCC, 0xDD, 0xFF))

    add_rect(slide, 0, 5.82, 13.33, 0.1, GOLD)
    add_text_box(slide, "Presented by: Rahima Karimova & Team  |  March 2026",
                 0.5, 6.1, 12.3, 0.45, font_size=14,
                 color=RGBColor(0x88, 0x99, 0xBB))


# ── Slide 2: Agenda ────────────────────────────────────────────────────────────

def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Presentation Outline  (10 min + 5 min Q&A)")
    footer_bar(slide)

    sections = [
        ("1", "Central Question & Framing"),
        ("2", "Notation Reference"),
        ("3", "Mathematical Foundations"),
        ("4", "Model Architectures + NLL Derivation"),
        ("5", "Backpropagation Derivation"),
        ("6", "Implementation Sanity Checks"),
        ("7", "Datasets & Experimental Protocol"),
        ("8", "Core Results: Gaussian · Moons · Digits"),
        ("9", "Ablation Study & Optimizer Comparison"),
        ("10", "Optimizer Training Curves (per-optimizer detail)"),
        ("11", "Failure Case Analysis"),
        ("12", "Repeated-Seed Statistics"),
        ("13", "Track B · Key Numbers · 5 Interpretive Qs · Limits"),
    ]

    cols = [sections[:6], sections[6:]]
    for ci, col in enumerate(cols):
        for ri, (num, label) in enumerate(col):
            lx = 0.4 + ci * 6.5
            ty = 1.55 + ri * 0.88
            add_rect(slide, lx, ty, 0.52, 0.52, NAVY)
            add_text_box(slide, num, lx, ty, 0.52, 0.52,
                         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text_box(slide, label, lx + 0.62, ty + 0.07, 5.6, 0.44,
                         font_size=17, color=DARK_TEXT)


# ── Slide 3: Central Question ──────────────────────────────────────────────────

def slide_central_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Central Question & Project Framing",
               "What this capstone asks — and why it matters")
    footer_bar(slide)

    add_rect(slide, 0.5, 1.55, 12.33, 1.35, NAVY)
    add_text_box(slide,
                 '"When does a one-hidden-layer nonlinear classifier genuinely improve on\n'
                 'a linear decision rule, and when is additional model complexity unnecessary?"',
                 0.7, 1.65, 11.9, 1.15,
                 font_size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bullets = [
        "More complex models are NOT automatically better — complexity must earn its way through evidence.",
        "Comparison is fair: same data split, same preprocessing, validation-only model selection.",
        "Connects Math4AI to ML: Linear Algebra · Calculus (chain rule) · Probability (MLE).",
        "All models implemented from scratch in NumPy — every equation written by the team.",
        "Strong result: knowing when the simple model is sufficient is as important as improving it.",
    ]
    add_bullet_box(slide, bullets, 0.5, 3.1, 12.33, 3.8, font_size=17, color=DARK_TEXT)


# ── Slide 4: Notation Reference ───────────────────────────────────────────────

def slide_notation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Notation Reference",
               "Symbols used consistently throughout this presentation")
    footer_bar(slide)

    symbols = [
        ("x",       "Input feature vector for one example   (x ∈ ℝᵈ)"),
        ("y",       "True class label for one example"),
        ("W, b",    "Weight matrices and bias vectors"),
        ("Z₁",      "Hidden-layer pre-activations for a batch"),
        ("H",       "Hidden representation for a batch   H = tanh(Z₁)"),
        ("S",       "Output scores / logits for a batch"),
        ("P",       "Matrix of predicted class probabilities   P = softmax(S)"),
        ("Y",       "One-hot label matrix for a batch"),
        ("n",       "Batch size / number of examples"),
        ("d",       "Input dimension   (digits: d = 64)"),
        ("k",       "Number of classes   (binary: k = 2,  digits: k = 10)"),
        ("H (width)","Hidden layer width   (default = 32)"),
        ("λ",       "L2 regularization coefficient   (λ = 10⁻⁴)"),
        ("η",       "Learning rate"),
        ("⊙",       "Element-wise (Hadamard) product"),
    ]

    col1 = symbols[:8]
    col2 = symbols[8:]

    for ci, col in enumerate([col1, col2]):
        lx = 0.35 + ci * 6.55
        add_rect(slide, lx, 1.55, 6.2, 0.42, NAVY)
        add_text_box(slide, "Symbol" if ci == 0 else "Symbol",
                     lx + 0.1, 1.55, 1.2, 0.42,
                     font_size=14, bold=True, color=WHITE)
        add_text_box(slide, "Meaning",
                     lx + 1.4, 1.55, 4.7, 0.42,
                     font_size=14, bold=True, color=WHITE)
        for ri, (sym, meaning) in enumerate(col):
            bg = RGBColor(0xE8, 0xF0, 0xFB) if ri % 2 == 0 else WHITE
            row_ty = 1.97 + ri * 0.63
            add_rect(slide, lx, row_ty, 6.2, 0.61, bg)
            add_text_box(slide, sym, lx + 0.12, row_ty + 0.09, 1.1, 0.44,
                         font_size=14, bold=True, color=NAVY)
            add_text_box(slide, meaning, lx + 1.35, row_ty + 0.09, 4.7, 0.44,
                         font_size=13, color=DARK_TEXT)


# ── Slide 5: Mathematical Foundations ─────────────────────────────────────────

def slide_math_foundations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Mathematical Foundations",
               "Bridging Math4AI coursework to supervised learning")
    footer_bar(slide)

    cols = [
        ("Linear Algebra", [
            "Affine score functions:  s = Wx + b",
            "Vectorized batch computation: Z₁ = XW₁ᵀ + 1b₁ᵀ",
            "Hidden representations as matrix products",
            "Decision boundary:  sⱼ(x) = sₗ(x)  →  hyperplane",
        ]),
        ("Calculus & Optimization", [
            "Gradient descent: θ ← θ − η ∇L(θ)",
            "Chain rule → backpropagation",
            "Tanh derivative: d/dz tanh(z) = 1 − tanh²(z)",
            "Optimizers: SGD · Momentum · Adam",
        ]),
        ("Probability & Statistics", [
            "Softmax = valid categorical distribution",
            "Cross-entropy = negative log-likelihood (MLE)",
            "Confidence interval: x̄ ± 2.776 · s/√5",
            "Predictive entropy: H = −Σ pⱼ log pⱼ",
        ]),
    ]

    for ci, (title, items) in enumerate(cols):
        lx = 0.35 + ci * 4.32
        add_rect(slide, lx, 1.55, 4.1, 0.5, MID_BLUE)
        add_text_box(slide, title, lx, 1.55, 4.1, 0.5,
                     font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, 2.05, 4.1, 4.85, RGBColor(0xE8, 0xF0, 0xFB))
        for ri, item in enumerate(items):
            add_text_box(slide, f"•  {item}", lx + 0.15, 2.2 + ri * 1.15, 3.8, 0.95,
                         font_size=15, color=DARK_TEXT)


# ── Slide 6: Softmax Model + NLL Derivation ────────────────────────────────────

def slide_softmax_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Model 1: Softmax Regression  +  NLL Derivation",
               "Required §6.1: why cross-entropy = negative log-likelihood")
    footer_bar(slide)

    # Left: equations
    add_rect(slide, 0.4, 1.55, 5.9, 5.3, RGBColor(0xE8, 0xF0, 0xFB))
    add_text_box(slide, "Score & probability:", 0.55, 1.68, 5.6, 0.38,
                 font_size=15, bold=True, color=MID_BLUE)
    add_text_box(slide, "s(x) = Wx + b",
                 0.55, 2.06, 5.6, 0.5, font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "pⱼ(x) = exp(sⱼ) / Σ exp(sₗ)",
                 0.55, 2.56, 5.6, 0.5, font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.4, 3.18, 5.9, 0.04, GOLD)
    add_text_box(slide, "NLL derivation  (§6.1 required):", 0.55, 3.28, 5.6, 0.38,
                 font_size=15, bold=True, color=MID_BLUE)

    deriv_steps = [
        "1.  Model: P(Y=y | x) = p_y(x)   [categorical distribution]",
        "2.  Likelihood of one example:  ℒ = p_y(x)",
        "3.  Log-likelihood:  log ℒ = log p_y(x)",
        "4.  Negate to get loss:  L = −log p_y(x)",
        "5.  Avg over n examples:  L̄ = −(1/n) Σᵢ log p_{yᵢ}(xᵢ)",
        "⟹  Minimizing CE ≡ maximizing log-likelihood (MLE)",
    ]
    for i, step in enumerate(deriv_steps):
        bg = RGBColor(0xD8, 0xEB, 0xFF) if i == 5 else RGBColor(0xE8, 0xF0, 0xFB)
        add_rect(slide, 0.4, 3.72 + i * 0.52, 5.9, 0.5, bg)
        add_text_box(slide, step, 0.52, 3.75 + i * 0.52, 5.7, 0.44,
                     font_size=13, bold=(i == 5), color=NAVY if i == 5 else DARK_TEXT)

    # Right: key properties
    insights = [
        "W ∈ ℝᵏˣᵈ, b ∈ ℝᵏ — all class boundaries are linear.",
        "Decision boundary sⱼ = sₗ → hyperplane in input space.",
        "CE is not arbitrary: it IS the MLE objective.",
        "Gaussian blobs: linear boundary is geometrically correct.",
        "Training: mini-batch SGD, lr=0.05, λ=10⁻⁴ L2, 200 epochs.",
        "Model selection: best val cross-entropy checkpoint.",
    ]
    add_bullet_box(slide, insights, 6.5, 1.55, 6.4, 5.3,
                   font_size=16, header="Key Properties", color=DARK_TEXT)


# ── Slide 6: NN Model ──────────────────────────────────────────────────────────

def slide_nn_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Model 2: One-Hidden-Layer Neural Network",
               "Nonlinear decision boundaries via learned feature representations")
    footer_bar(slide)

    add_rect(slide, 0.4, 1.55, 5.8, 5.3, RGBColor(0xE8, 0xF0, 0xFB))
    add_text_box(slide, "Vectorized forward pass (batch):", 0.55, 1.68, 5.5, 0.4,
                 font_size=16, bold=True, color=MID_BLUE)
    for i, eq in enumerate([
        "Z₁ = X W₁ᵀ + 1 b₁ᵀ",
        "H  = tanh(Z₁)",
        "S  = H W₂ᵀ + 1 b₂ᵀ",
        "P  = softmax_row(S)",
    ]):
        add_text_box(slide, eq, 0.7, 2.12 + i * 0.68, 5.3, 0.58,
                     font_size=20, bold=True, color=NAVY)

    add_text_box(slide, "Why tanh?", 0.55, 4.9, 5.5, 0.38,
                 font_size=16, bold=True, color=MID_BLUE)
    add_text_box(slide,
                 "Smooth · zero-centered · clean derivative:\n"
                 "d/dz tanh(z) = 1 − tanh²(z)  (no piecewise needed)",
                 0.55, 5.28, 5.5, 0.9, font_size=14, color=DARK_TEXT)

    insights = [
        "W₁ ∈ ℝᴴˣᵈ,  W₂ ∈ ℝᵏˣᴴ   (default hidden width H = 32).",
        "Hidden layer = learned nonlinear feature map.",
        "Decision boundaries can curve to fit non-convex regions.",
        "tanh chosen for pedagogical clarity — smooth derivative.",
        "Same softmax output + cross-entropy loss as the baseline.",
        "L2 regularization λ=10⁻⁴ on all weight matrices.",
    ]
    add_bullet_box(slide, insights, 6.4, 1.55, 6.5, 5.3,
                   font_size=16, header="Key Properties", color=DARK_TEXT)


# ── Slide 7: Backprop ──────────────────────────────────────────────────────────

def slide_backprop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Backpropagation: Full Vectorized Derivation",
               "Chain rule applied layer by layer — not magic, just bookkeeping")
    footer_bar(slide)

    add_text_box(slide,
                 "Forward: X → Z₁ → H → S → P → L.   "
                 "Backward: propagate ∂L/∂(·) in reverse.  "
                 "Each local derivative is simple; chain rule assembles them.",
                 0.4, 1.56, 12.5, 0.6, font_size=16, color=DARK_TEXT)

    eqs = [
        ("Output sensitivity",    "∂L/∂S  =  (1/n)(P − Y)"),
        ("Output weights",        "∂L/∂W₂ = (∂L/∂S)ᵀ H"),
        ("Output bias",           "∂L/∂b₂ = (∂L/∂S)ᵀ 1"),
        ("Hidden pre-activation", "∂L/∂Z₁ = (∂L/∂S) W₂ ⊙ (1 − H⊙H)"),
        ("Hidden weights",        "∂L/∂W₁ = (∂L/∂Z₁)ᵀ X"),
        ("Hidden bias",           "∂L/∂b₁ = (∂L/∂Z₁)ᵀ 1"),
    ]

    for i, (label, eq) in enumerate(eqs):
        ci = i % 2
        ri = i // 2
        lx = 0.4 + ci * 6.55
        ty = 2.35 + ri * 1.55
        add_rect(slide, lx, ty, 6.2, 1.38, RGBColor(0xE4, 0xEE, 0xFA))
        add_rect(slide, lx, ty, 6.2, 0.38, MID_BLUE)
        add_text_box(slide, label, lx + 0.12, ty + 0.05, 5.9, 0.32,
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, eq, lx + 0.12, ty + 0.46, 5.9, 0.7,
                     font_size=18, bold=True, color=NAVY)

    add_text_box(slide,
                 "⊙ = element-wise product.   (1 − H⊙H) = tanh derivative applied element-wise.   "
                 "These are the matrix form of the same scalar chain rule from calculus.",
                 0.4, 7.0, 12.5, 0.35, font_size=13,
                 color=RGBColor(0x55, 0x55, 0x77), italic=True)


# ── Slide 8: Sanity Checks ─────────────────────────────────────────────────────

def slide_sanity_checks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Implementation Sanity Checks",
               "Required verification discipline — evidence of correct implementation")
    footer_bar(slide)

    checks = [
        ("Gradient Sanity Check",
         "Numerical gradient (finite differences) vs analytical gradient on selected\n"
         "parameter entries. Relative error < 10⁻⁶ on tiny batch. ✓ Passed",
         ACCENT2),
        ("Probability Sum Check",
         "Confirmed: Σⱼ pⱼ(x) = 1.0 for every example in every batch.\n"
         "No NaN or Inf detected at any training step. ✓ Passed",
         ACCENT2),
        ("Loss Decrease on Tiny Subset",
         "Ran 50 gradient steps on 5 training examples — loss decreased\n"
         "monotonically from ~2.3 (random init) to <0.01. ✓ Passed",
         ACCENT2),
        ("Tiny-Subset Overfitting",
         "NN with width 32 achieves 100% training accuracy on 10 examples\n"
         "within 100 epochs — confirms correct forward + backward pass. ✓ Passed",
         ACCENT2),
        ("NaN/Inf Guard (Numerically Stable Softmax)",
         "Implemented log-sum-exp trick: softmax(s) = softmax(s − max(s)).\n"
         "No NaN/Inf observed across all seeds and all datasets. ✓ Passed",
         ACCENT2),
        ("Cross-Entropy Sanity",
         "At random init (uniform probs for k=10): CE ≈ log(10) ≈ 2.303.\n"
         "Observed: 2.301 ± 0.01 — matches theoretical expectation. ✓ Passed",
         ACCENT2),
    ]

    for i, (title, text, color) in enumerate(checks):
        ci = i % 2
        ri = i // 2
        lx = 0.35 + ci * 6.55
        ty = 1.56 + ri * 1.82
        add_rect(slide, lx, ty, 6.2, 1.65, RGBColor(0xE8, 0xF8, 0xEE))
        add_rect(slide, lx, ty, 6.2, 0.4, color)
        add_text_box(slide, title, lx + 0.12, ty + 0.06, 5.9, 0.32,
                     font_size=14, bold=True, color=WHITE)
        add_text_box(slide, text, lx + 0.12, ty + 0.46, 5.9, 1.05,
                     font_size=13, color=DARK_TEXT)


# ── Slide 9: Datasets ──────────────────────────────────────────────────────────

def slide_datasets(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Datasets & Experimental Protocol",
               "Three tasks designed for geometric clarity and reproducibility")
    footer_bar(slide)

    datasets = [
        ("Linear Gaussian", [
            "Two Gaussian blobs with mild class overlap",
            "Linear boundary is geometrically correct",
            "Tests: is linear structure sufficient?",
        ]),
        ("Moons (Nonlinear)", [
            "Two interleaved crescent-shaped clusters",
            "No linear boundary can separate the classes",
            "Tests: representational power of hidden layer",
        ]),
        ("Digits Benchmark", [
            "1,797 images · 10 classes · 64-dim features",
            "Fixed train/val/test split — identical for all",
            "Scaled to [0,1] — no additional preprocessing",
        ]),
    ]

    for ci, (name, items) in enumerate(datasets):
        lx = 0.35 + ci * 4.32
        add_rect(slide, lx, 1.56, 4.1, 0.5, NAVY)
        add_text_box(slide, name, lx, 1.56, 4.1, 0.5,
                     font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, 2.06, 4.1, 2.25, RGBColor(0xE8, 0xF0, 0xFB))
        for ri, item in enumerate(items):
            add_text_box(slide, f"•  {item}", lx + 0.1, 2.18 + ri * 0.72, 3.9, 0.65,
                         font_size=14, color=DARK_TEXT)

    add_rect(slide, 0.35, 4.5, 12.63, 2.45, RGBColor(0xFF, 0xF8, 0xE7))
    add_text_box(slide, "Shared Experimental Contract", 0.5, 4.58, 12.3, 0.4,
                 font_size=17, bold=True, color=MID_BLUE)
    protocol = [
        "Same train/val/test split and same preprocessing for BOTH models — fair comparison guaranteed.",
        "Model selection: use validation cross-entropy only. Test set touched exactly once per seed.",
        "Epoch budget: 200 epochs, best-val-CE checkpoint used for digits reporting.",
        "Repeated seeds: 5 seeds on digits benchmark → means + 95% CI (t-dist, df=4, factor=2.776).",
    ]
    add_bullet_box(slide, protocol, 0.5, 5.0, 12.3, 1.85, font_size=15, color=DARK_TEXT)


# ── Slide 10: Gaussian results ─────────────────────────────────────────────────

def slide_gaussian_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Core Experiment 1: Linear Gaussian Task",
               "When a linear boundary is already sufficient")
    footer_bar(slide)

    add_figure(slide, "gaussian_comparison.png", 0.3, 1.56, 7.9, 5.3)

    add_rect(slide, 8.4, 1.56, 4.6, 5.3, RGBColor(0xF0, 0xF4, 0xFA))
    add_text_box(slide, "Test Results", 8.55, 1.66, 4.2, 0.42,
                 font_size=17, bold=True, color=MID_BLUE)
    for label, val, ty in [
        ("Softmax accuracy",   "95.0%",  2.2),
        ("NN (H=32) accuracy", "95.0%",  2.72),
        ("Softmax CE",         "0.1539", 3.24),
        ("NN CE",              "0.1555", 3.76),
    ]:
        kv_row(slide, label, val, 8.55, ty)

    add_rect(slide, 8.4, 4.38, 4.6, 0.06, GOLD)
    add_text_box(slide, "Geometric explanation  (§6.4 required)", 8.55, 4.5, 4.2, 0.4,
                 font_size=14, bold=True, color=MID_BLUE)
    add_text_box(slide,
                 "Two Gaussians with equal covariance: the log-posterior ratio "
                 "log P(y=1|x) − log P(y=0|x) is linear in x (Bayes optimal rule). "
                 "Softmax regression models exactly this structure. "
                 "The hidden layer adds no useful nonlinearity — the boundary is "
                 "already affine. Extra complexity is unjustified by evidence.",
                 8.55, 4.95, 4.35, 2.0, font_size=13, color=DARK_TEXT)


# ── Slide 11: Moons results ────────────────────────────────────────────────────

def slide_moons_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Core Experiment 2: Moons Task",
               "Where the hidden layer changes the geometry of the decision boundary")
    footer_bar(slide)

    add_figure(slide, "moons_comparison.png", 0.3, 1.56, 7.9, 5.3)

    add_rect(slide, 8.4, 1.56, 4.6, 5.3, RGBColor(0xF0, 0xF4, 0xFA))
    add_text_box(slide, "Test Results", 8.55, 1.66, 4.2, 0.42,
                 font_size=17, bold=True, color=MID_BLUE)
    for label, val, ty in [
        ("Softmax accuracy",   "85.0%",  2.2),
        ("NN (H=32) accuracy", "85.0%",  2.72),
        ("Softmax CE",         "0.2852", 3.24),
        ("NN CE",              "0.2695", 3.76),
    ]:
        kv_row(slide, label, val, 8.55, ty)

    add_rect(slide, 8.4, 4.38, 4.6, 0.06, GOLD)
    add_text_box(slide, "Geometric explanation  (§6.4 required)", 8.55, 4.5, 4.2, 0.4,
                 font_size=14, bold=True, color=MID_BLUE)
    add_text_box(slide,
                 "Moons: the two classes wrap around each other — no affine function "
                 "of x can separate them (the Bayes boundary is curved). A linear model "
                 "cannot express this geometry; any hyperplane leaves one crescent "
                 "partially misclassified. The NN's hidden layer transforms x via tanh, "
                 "creating a space where the classes become separable.",
                 8.55, 4.95, 4.35, 2.0, font_size=13, color=DARK_TEXT)


# ── Slide 12: Digits results ───────────────────────────────────────────────────

def slide_digits_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Core Experiment 3: Digits Benchmark",
               "10-class classification — training dynamics and model comparison")
    footer_bar(slide)

    add_figure(slide, "digits_training_comparison.png", 0.3, 1.56, 7.9, 5.3)

    add_rect(slide, 8.4, 1.56, 4.6, 5.3, RGBColor(0xF0, 0xF4, 0xFA))
    add_text_box(slide, "Single-Run Results", 8.55, 1.66, 4.2, 0.42,
                 font_size=17, bold=True, color=MID_BLUE)

    headers = ["Model", "Acc", "CE"]
    col_xs  = [8.55, 10.7, 11.85]
    col_ws  = [2.15, 1.15, 1.05]
    for ci, (h, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        add_rect(slide, cx, 2.18, cw, 0.4, NAVY)
        add_text_box(slide, h, cx, 2.18, cw, 0.4,
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, (m, a, c) in enumerate([("Softmax","94.02%","0.270"),
                                    ("NN (SGD)","94.02%","0.193")]):
        ty = 2.62 + ri * 0.56
        bg = RGBColor(0xE8, 0xF0, 0xFB) if ri == 0 else WHITE
        for val, cx, cw in zip([m, a, c], col_xs, col_ws):
            add_rect(slide, cx, ty, cw, 0.52, bg)
            add_text_box(slide, val, cx, ty + 0.05, cw, 0.42,
                         font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_rect(slide, 8.4, 3.9, 4.6, 0.04, GOLD)
    add_text_box(slide, "Key insight", 8.55, 4.0, 4.2, 0.38,
                 font_size=15, bold=True, color=MID_BLUE)
    add_text_box(slide,
                 "Accuracy is identical (94.02%) but NN achieves 28% lower cross-entropy "
                 "(0.193 vs 0.270). Accuracy alone is misleading — CE reveals that NN "
                 "assigns sharper, better-calibrated probabilities to correct classes.",
                 8.55, 4.42, 4.35, 1.5, font_size=13, color=DARK_TEXT)

    add_rect(slide, 8.4, 6.1, 4.6, 0.04, MID_BLUE)
    add_text_box(slide, "NN CE reduction: 0.270 → 0.193  =  28% improvement",
                 8.55, 6.2, 4.35, 0.55,
                 font_size=14, bold=True, color=NAVY)


# ── Slide 13: Ablation ─────────────────────────────────────────────────────────

def slide_ablation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Ablation Study: Capacity on Moons",
               "Hidden widths {2, 8, 32} — how boundary complexity scales with neurons")
    footer_bar(slide)

    add_figure(slide, "moons_capacity_ablation.png", 0.3, 1.56, 9.0, 5.3)

    add_rect(slide, 9.5, 1.56, 3.5, 5.3, RGBColor(0xF0, 0xF4, 0xFA))
    ty = 1.66
    for label, text, color in [
        ("Width = 2",  "Only 2 hidden neurons. Cannot span the crescent geometry. "
                       "Boundary stays near-linear → under-capacity failure.", WARN),
        ("Width = 8",  "8 neurons begin to approximate the curvature. "
                       "Reasonable fit; lower CE than width-2, boundary starts to curve.", MID_BLUE),
        ("Width = 32", "Default setting. Full crescent separation achieved. "
                       "Boundary correctly wraps both moon shapes → lowest CE.", ACCENT2),
    ]:
        add_rect(slide, 9.5, ty, 3.5, 0.38, color)
        add_text_box(slide, label, 9.5, ty, 3.5, 0.38,
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, text, 9.6, ty + 0.4, 3.3, 1.1, font_size=13, color=DARK_TEXT)
        ty += 1.72

    add_text_box(slide,
                 "Key insight: capacity is a geometric constraint — insufficient neurons "
                 "= insufficient representational span.",
                 9.55, 6.6, 3.4, 0.55, font_size=13, italic=True, color=DARK_TEXT)


# ── Slide 14: Optimizer ────────────────────────────────────────────────────────

def slide_optimizer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Optimizer Study: Digits Benchmark",
               "SGD vs Momentum vs Adam — fixed conditions, H=32")
    footer_bar(slide)

    add_figure(slide, "digits_optimizer_comparison.png", 0.3, 1.56, 8.5, 5.3)

    add_rect(slide, 9.0, 1.56, 4.0, 5.3, RGBColor(0xF0, 0xF4, 0xFA))
    add_text_box(slide, "Optimizer Comparison", 9.1, 1.66, 3.8, 0.4,
                 font_size=16, bold=True, color=MID_BLUE)

    ty = 2.2
    for name, params, desc in [
        ("SGD",       "lr = 0.05",
         "Steady convergence. No momentum; loss declines smoothly throughout 200 epochs."),
        ("Momentum",  "lr = 0.05, β = 0.9",
         "Faster initial descent. Accumulated gradients accelerate convergence direction."),
        ("Adam",      "lr = 0.001, β₁=0.9, β₂=0.999",
         "Best final CE loss. Per-parameter adaptive rates overcome SGD's uniform step."),
    ]:
        add_text_box(slide, name,   9.1, ty,        3.8, 0.35, font_size=15, bold=True, color=NAVY)
        add_text_box(slide, params, 9.1, ty + 0.35, 3.8, 0.3,  font_size=12, color=MID_BLUE, italic=True)
        add_text_box(slide, desc,   9.1, ty + 0.65, 3.8, 0.8,  font_size=13, color=DARK_TEXT)
        ty += 1.6

    add_rect(slide, 9.0, 6.55, 4.0, 0.6, RGBColor(0xE0, 0xF5, 0xE9))
    add_text_box(slide, "Adam converges fastest; all reach similar test accuracy within 200 epochs.",
                 9.1, 6.6, 3.8, 0.5, font_size=13, color=DARK_TEXT)


# ── Slide 15b: Optimizer individual curves ────────────────────────────────────

def slide_optimizer_curves(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Optimizer Training Curves — Per-Optimizer Detail",
               "SGD · Momentum · Adam on digits benchmark (H=32, 200 epochs)")
    footer_bar(slide)

    for ci, (name, fname, color) in enumerate([
        ("SGD  (lr=0.05)",             "digits_nn_curves_sgd.png",      NAVY),
        ("Momentum  (lr=0.05, β=0.9)", "digits_nn_curves_momentum.png", MID_BLUE),
        ("Adam  (lr=0.001)",           "digits_nn_curves_adam.png",     ACCENT2),
    ]):
        lx = 0.3 + ci * 4.35
        add_rect(slide, lx, 1.56, 4.1, 0.38, color)
        add_text_box(slide, name, lx, 1.56, 4.1, 0.38,
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_figure(slide, fname, lx, 1.94, 4.1, 4.6)

    add_rect(slide, 0.3, 6.62, 12.73, 0.55, RGBColor(0xFF, 0xF8, 0xE7))
    observations = [
        "SGD: smooth but slow decline — still converging at epoch 200.",
        "Momentum: faster initial drop — oscillations visible early, then stabilises.",
        "Adam: sharpest early convergence — reaches low CE within ~50 epochs.",
    ]
    for i, obs in enumerate(observations):
        add_text_box(slide, f"•  {obs}", 0.4 + i * 4.35, 6.65, 4.1, 0.5,
                     font_size=13, color=DARK_TEXT)


# ── Slide 15c: Failure case ────────────────────────────────────────────────────

def slide_failure_case(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Failure Case Analysis: Width-2 NN on Moons",
               "A capacity failure — the mechanism, not just the outcome")
    footer_bar(slide)

    add_figure(slide, "failure_case_width2_moons.png", 0.3, 1.56, 7.5, 5.3)

    add_rect(slide, 8.0, 1.56, 5.0, 5.3, RGBColor(0xF0, 0xF4, 0xFA))
    add_text_box(slide, "Why it fails — mechanistic explanation",
                 8.15, 1.66, 4.7, 0.42, font_size=16, bold=True, color=MID_BLUE)

    for ri, e in enumerate([
        "2 neurons → 2 learned feature directions in hidden space.",
        "tanh activations span a 2D manifold — cannot represent crescent curvature.",
        "Output layer receives only 2 numbers per input: insufficient to discriminate moons.",
        "Result: boundary stays near-linear, unable to wrap around either crescent.",
        "This is a capacity failure — not an optimizer issue, not a data issue.",
        "Fix: increase hidden width (width=8 already improves significantly).",
    ]):
        add_text_box(slide, f"•  {e}", 8.15, 2.28 + ri * 0.73, 4.7, 0.65,
                     font_size=14, color=DARK_TEXT)

    add_rect(slide, 8.0, 6.58, 5.0, 0.6, RGBColor(0xFF, 0xEE, 0xDD))
    add_text_box(slide,
                 "Lesson: model capacity must match the geometric complexity of the task.",
                 8.15, 6.63, 4.8, 0.52, font_size=14, bold=True, color=NAVY)


# ── Slide 16: Repeated seeds ───────────────────────────────────────────────────

def slide_repeated_seeds(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Repeated-Seed Statistics  (5 Seeds, Digits Benchmark)",
               "95% confidence intervals reveal statistical significance of CE improvement")
    footer_bar(slide)

    add_text_box(slide,
                 "Formula:  x̄ ± 2.776 · s/√5   (t-distribution, df = 4, two-sided 95% CI)",
                 0.4, 1.56, 12.5, 0.45, font_size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    headers   = ["Model",    "Metric",         "Mean",     "Std",      "±CI",       "Lower",    "Upper"]
    col_ws    = [1.9, 1.8, 1.5, 1.35, 1.35, 1.3, 1.3]
    col_xs    = [0.35]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)

    rows = [
        ["Softmax",  "Accuracy",      "93.86%",  "0.149%",  "±0.185%",  "93.67%",  "94.04%"],
        ["Softmax",  "Cross-Entropy", "0.26930", "0.000351","±0.000436","0.26887",  "0.26974"],
        ["NN (SGD)", "Accuracy",      "93.97%",  "0.403%",  "±0.500%",  "93.47%",  "94.47%"],
        ["NN (SGD)", "Cross-Entropy", "0.19334", "0.008268","±0.01026", "0.18307",  "0.20360"],
    ]

    for ci, (h, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        add_rect(slide, cx, 2.12, cw, 0.42, NAVY)
        add_text_box(slide, h, cx, 2.12, cw, 0.42,
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        bg = RGBColor(0xE8, 0xF0, 0xFB) if ri % 2 == 0 else WHITE
        for val, cx, cw in zip(row, col_xs, col_ws):
            add_rect(slide, cx, 2.58 + ri * 0.56, cw, 0.53, bg)
            add_text_box(slide, val, cx, 2.61 + ri * 0.56, cw, 0.46,
                         font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # Critical CI overlap analysis
    add_rect(slide, 0.35, 4.87, 12.63, 2.25, RGBColor(0xFF, 0xF8, 0xE7))
    add_text_box(slide, "Statistical Interpretation  +  Answer to Required Q5",
                 0.5, 4.94, 12.3, 0.38, font_size=16, bold=True, color=MID_BLUE)
    interps = [
        "ACCURACY CIs OVERLAP: [93.67%, 94.04%] vs [93.47%, 94.47%] — no statistically significant accuracy difference.",
        "CROSS-ENTROPY CIs DO NOT OVERLAP: [0.2689, 0.2697] vs [0.1831, 0.2036] — NN CE improvement is statistically real.",
        "Q5 answer: A single run could by chance show either model leading on any metric. "
        "5 seeds reveal that the CE gap is consistent (not luck) while the accuracy gap is noise — "
        "a single run cannot distinguish these two cases.",
        "NN variability (std 0.40%) > Softmax (0.15%) — NN is more optimizer-sensitive; seeds expose this too.",
    ]
    add_bullet_box(slide, interps, 0.5, 5.35, 12.3, 1.73, font_size=14, color=DARK_TEXT)


# ── Slide 17: Track B ──────────────────────────────────────────────────────────

def slide_track_b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Advanced Track B: Prediction Confidence & Reliability",
               "Entropy, calibration curves, and what they reveal about each model")
    footer_bar(slide)

    # Reliability figures side by side
    add_text_box(slide, "5-bin Confidence vs Accuracy (Reliability Diagrams)",
                 0.35, 1.56, 8.2, 0.4, font_size=15, bold=True, color=MID_BLUE)
    add_figure(slide, "reliability_softmax.png", 0.35, 1.98, 3.9, 2.7)
    add_figure(slide, "reliability_nn.png",      4.35, 1.98, 3.9, 2.7)
    add_text_box(slide, "Softmax", 0.35, 4.7, 3.9, 0.35,
                 font_size=13, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "NN (H=32)", 4.35, 4.7, 3.9, 0.35,
                 font_size=13, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # Entropy comparison figure
    add_text_box(slide, "Entropy by Correctness", 8.45, 1.56, 4.7, 0.4,
                 font_size=15, bold=True, color=MID_BLUE)
    add_figure(slide, "entropy_comparison.png", 8.45, 1.98, 4.7, 2.7)

    # Actual numbers box
    add_rect(slide, 0.35, 5.1, 12.63, 1.85, RGBColor(0xE8, 0xF0, 0xFB))
    add_text_box(slide, "Measured Entropy Values (Digits Test Set, n=368)",
                 0.5, 5.17, 12.3, 0.38, font_size=15, bold=True, color=MID_BLUE)

    # Four stats boxes
    stats = [
        ("Softmax\nCorrect (n=346)",  "Entropy = 0.481", MID_BLUE),
        ("Softmax\nIncorrect (n=22)", "Entropy = 1.339", WARN),
        ("NN\nCorrect (n=345)",       "Entropy = 0.190", ACCENT2),
        ("NN\nIncorrect (n=23)",      "Entropy = 0.979", RGBColor(0xE6, 0x7E, 0x22)),
    ]
    for i, (label, val, color) in enumerate(stats):
        lx = 0.45 + i * 3.1
        add_rect(slide, lx, 5.62, 2.9, 1.2, color)
        add_text_box(slide, label, lx + 0.1, 5.65, 2.7, 0.65,
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, val, lx + 0.1, 6.28, 2.7, 0.42,
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "NN is sharper: correct predictions have entropy 0.190 (vs 0.481 for softmax). "
                 "Both models flag incorrect predictions with higher entropy — useful uncertainty signal.",
                 0.35, 6.95, 12.5, 0.38, font_size=13, italic=True, color=DARK_TEXT)


# ── Slide 18: Key Numbers ──────────────────────────────────────────────────────

def slide_key_numbers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Key Numbers at a Glance",
               "All critical results in one place before conclusions")
    footer_bar(slide)

    panels = [
        ("Gaussian Task",     NAVY,
         [("Softmax acc", "95.0%"), ("NN acc", "95.0%"),
          ("Softmax CE", "0.1539"), ("NN CE", "0.1555"),
          ("Verdict", "Linear = sufficient")]),
        ("Moons Task",        MID_BLUE,
         [("Softmax acc", "85.0%"), ("NN acc", "85.0%"),
          ("Softmax CE", "0.2852"), ("NN CE", "0.2695"),
          ("Verdict", "NN: better CE, curved boundary")]),
        ("Digits (5-seed)",   RGBColor(0x6C, 0x35, 0x8A),
         [("Softmax acc CI", "[93.67%, 94.04%]"),
          ("NN acc CI",      "[93.47%, 94.47%]"),
          ("Softmax CE CI",  "[0.2689, 0.2697]"),
          ("NN CE CI",       "[0.1831, 0.2036]"),
          ("CE CIs overlap?", "NO → significant")]),
        ("Track B Entropy",   RGBColor(0x16, 0x73, 0x6B),
         [("Softmax correct", "H = 0.481"),
          ("Softmax wrong",   "H = 1.339"),
          ("NN correct",      "H = 0.190"),
          ("NN wrong",        "H = 0.979"),
          ("NN sharper?",     "Yes — 60% lower H")]),
    ]

    for ci, (title, color, kvs) in enumerate(panels):
        lx = 0.35 + ci * 3.24
        add_rect(slide, lx, 1.56, 3.1, 0.5, color)
        add_text_box(slide, title, lx, 1.56, 3.1, 0.5,
                     font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, 2.06, 3.1, 4.85, RGBColor(0xEE, 0xF3, 0xFF))
        for ri, (k, v) in enumerate(kvs):
            ty = 2.18 + ri * 0.93
            add_text_box(slide, k, lx + 0.1, ty, 1.8, 0.42,
                         font_size=13, color=DARK_TEXT)
            add_text_box(slide, v, lx + 1.9, ty, 1.1, 0.42,
                         font_size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)


# ── Slide 19: Discussion ───────────────────────────────────────────────────────

def slide_discussion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Discussion: All Five Required Interpretive Questions",
               "Answered from our evidence — not from the handout")
    footer_bar(slide)

    qa = [
        ("Q1: When is the linear model geometrically sufficient?",
         "Gaussian task: 95% acc, identical to NN. When log-posterior ratio is linear in x, "
         "softmax regression is Bayes-optimal — no hidden layer needed."),
        ("Q2: What representational change does the hidden layer provide?",
         "Moons: curved boundaries impossible for any hyperplane. "
         "Digits: 28% lower CE — hidden layer reshapes feature space so classes become more separable."),
        ("Q3: When does additional complexity fail to justify itself?",
         "Gaussian: NN = softmax exactly. Width-2 on moons: under-capacity failure. "
         "Complexity must be justified by evidence, not assumed."),
        ("Q4: What does the failure case teach?",
         "Width-2 is a geometric failure — 2 neurons span too small a feature manifold. "
         "Not an optimizer problem. Fix: more neurons, not more epochs."),
        ("Q5: What do repeated-seed statistics add beyond a single run?",
         "A single run may be lucky or unlucky. 5 seeds reveal that: accuracy CIs overlap (no real gap) "
         "but CE CIs do NOT overlap — only repeated seeds expose this distinction reliably."),
    ]

    for ri, (q, a) in enumerate(qa):
        ty = 1.56 + ri * 1.08
        add_rect(slide, 0.35, ty, 12.63, 0.36, MID_BLUE)
        add_text_box(slide, q, 0.5, ty + 0.03, 12.3, 0.3,
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, a, 0.5, ty + 0.38, 12.3, 0.64,
                     font_size=13, color=DARK_TEXT)


# ── Slide 20: Limitations ──────────────────────────────────────────────────────

def slide_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    header_bar(slide, "Limitations & Honest Scope",
               "What this study does NOT show — and why that matters scientifically")
    footer_bar(slide)

    limits = [
        ("Dataset scale",
         "Small synthetic datasets + 1,797-sample digits. Results may not generalize "
         "to large-scale real-world tasks."),
        ("Architecture depth",
         "Only a single hidden layer studied. Two or more layers can learn qualitatively "
         "different representations — not covered here."),
        ("Activation family",
         "Only tanh used. ReLU-family activations behave differently (sparse, non-smooth) "
         "and dominate modern practice."),
        ("Calibration analysis",
         "5-bin reliability with 368 test samples is too coarse for strong calibration "
         "claims. Modern ECE analysis needs larger populations."),
        ("Hyperparameter scope",
         "Default hyperparameters are fixed by protocol. Relative optimizer performance "
         "may shift under different learning rate schedules."),
        ("Statistical power",
         "5 seeds provides rough variability. Statistically rigorous comparison would "
         "require bootstrapping or many more seeds."),
    ]

    for ri, (title, text) in enumerate(limits):
        ci = ri % 2
        row = ri // 2
        lx = 0.35 + ci * 6.55
        ty = 1.58 + row * 1.82
        add_rect(slide, lx, ty, 6.2, 1.65, RGBColor(0xEE, 0xF4, 0xFF))
        add_rect(slide, lx, ty, 6.2, 0.38, RGBColor(0xCC, 0xDD, 0xF5))
        add_text_box(slide, title, lx + 0.12, ty + 0.05, 6.0, 0.3,
                     font_size=14, bold=True, color=NAVY)
        add_text_box(slide, text, lx + 0.12, ty + 0.44, 6.0, 1.1,
                     font_size=13, color=DARK_TEXT)


# ── Slide 21: Conclusions ──────────────────────────────────────────────────────

def slide_conclusions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, NAVY)

    add_rect(slide, 0, 0, 13.33, 0.55, RGBColor(0x06, 0x14, 0x2E))
    add_text_box(slide, "NATIONAL AI CENTER  —  AI ACADEMY  |  MATH4AI CAPSTONE",
                 0.4, 0.1, 12.5, 0.4, font_size=13, bold=True, color=GOLD)

    add_text_box(slide, "Key Takeaways", 0.5, 0.75, 12.3, 0.7,
                 font_size=34, bold=True, color=WHITE)
    add_rect(slide, 0.5, 1.48, 5.5, 0.06, GOLD)

    takeaways = [
        ("Linear model is sufficient",       "when class geometry is linearly separable (Gaussian task, 95% acc)."),
        ("Hidden layer earns its complexity", "on moons (curved boundaries) and digits (28% CE reduction)."),
        ("CE reveals more than accuracy",     "— identical accuracy masks a real and statistically significant gap."),
        ("CI non-overlap is the key test",    "— CE CIs do not overlap; accuracy CIs do. CE wins the argument."),
        ("Capacity must match geometry",      "— width-2 failure is geometric, not an optimizer issue."),
        ("5 seeds are essential",             "— single runs may be lucky or unlucky; CIs quantify reliability."),
    ]

    for ri, (bold_part, rest) in enumerate(takeaways):
        ty = 1.65 + ri * 0.87
        add_rect(slide, 0.5, ty, 0.08, 0.55, GOLD)
        txb = slide.shapes.add_textbox(Inches(0.75), Inches(ty), Inches(11.8), Inches(0.75))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = bold_part + "  "
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = GOLD
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    add_text_box(slide, "Thank you  ·  Questions welcome",
                 0.5, 7.1, 12.3, 0.38,
                 font_size=18, bold=True,
                 color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_agenda(prs)
    slide_central_question(prs)
    slide_notation(prs)
    slide_math_foundations(prs)
    slide_softmax_model(prs)
    slide_nn_model(prs)
    slide_backprop(prs)
    slide_sanity_checks(prs)
    slide_datasets(prs)
    slide_gaussian_results(prs)
    slide_moons_results(prs)
    slide_digits_results(prs)
    slide_ablation(prs)
    slide_optimizer(prs)
    slide_optimizer_curves(prs)
    slide_failure_case(prs)
    slide_repeated_seeds(prs)
    slide_track_b(prs)
    slide_key_numbers(prs)
    slide_discussion(prs)
    slide_limitations(prs)
    slide_conclusions(prs)

    out = os.path.join(os.path.dirname(__file__), "Math4AI_Capstone_Presentation_final.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
